#!/usr/bin/env python3
"""
generate-avalere-pptx build engine (template-as-base).

Starts from the pinned Avalere PowerPoint template, adds slides by selecting
existing named layouts from its 132-layout catalogue, and injects content into
the layout placeholders. The theme (Avalere Health 2025), the Inter fonts, and
every branded layout are preserved because we never rebuild geometry, we only
fill placeholders on the template's own layouts.

Usage:
    python3 build_deck.py <spec.json> [--template <path>] [--out <path>]

Spec JSON shape:
{
  "output": "sample_deck.pptx",
  "title": "Deck title",
  "slides": [
    {
      "layout": "Cover - Pink - Custom",     # exact or fuzzy layout name
      "title": "Velorixa",
      "subtitle": "Brand strategy",
      "bullets": ["point one", "point two"], # optional
      "citations": ["623851f2"],              # raw source ids, rendered as [src:id]
      "notes": "speaker notes"                # optional
    }
  ]
}

Layout names are read from the template at build time (Appendix A note). Use
list-layouts to print the catalogue.
"""
import sys, os, json, argparse, zipfile, tempfile, shutil, re

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
DEFAULT_TEMPLATE = os.path.normpath(os.path.join(HERE, "..", "_assets", "Avalere_PPT_template.potx"))


def load_template(path):
    """Open a .potx/.pptx with python-pptx, normalising the template content type if needed."""
    try:
        return Presentation(path)
    except Exception:
        tmp = tempfile.mktemp(suffix=".pptx")
        with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item == "[Content_Types].xml":
                    data = data.replace(
                        b"presentationml.template.main+xml",
                        b"presentationml.presentation.main+xml",
                    )
                zout.writestr(item, data)
        return Presentation(tmp)


def clear_slides(prs):
    """Remove the template's demo slides so the deck starts clean, keeping masters/layouts/theme.

    Drops both the sldId entry and the presentation->slide relationship, so the orphaned
    slide parts are not reachable in the relationship graph and are not re-serialised
    (which would otherwise collide with new slideN.xml part names).
    """
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        sldIdLst.remove(sldId)
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass


def layout_map(prs):
    """name -> layout object, across every slide master."""
    m = {}
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            m.setdefault(layout.name.strip(), layout)
    return m


def _norm(s):
    """Lowercase and collapse dash variants/whitespace so specs can use plain hyphens."""
    s = s.replace("–", "-").replace("—", "-")
    return re.sub(r"\s+", " ", s).strip().lower()


def pick_layout(lmap, requested):
    """Exact match, else dash/case-insensitive, else first whose normalised name contains the request."""
    if requested in lmap:
        return lmap[requested]
    norm = {_norm(k): v for k, v in lmap.items()}
    r = _norm(requested)
    if r in norm:
        return norm[r]
    for k, v in norm.items():
        if r in k:
            return v
    raise SystemExit(f"Layout not found: {requested!r}. Run with 'list-layouts' to see options.")


def cite_suffix(citations):
    return (" " + " ".join(f"[src:{c}]" for c in citations)) if citations else ""


def fill_slide(slide, spec):
    cites = spec.get("citations", [])
    title_text = spec.get("title", "")
    subtitle = spec.get("subtitle")
    bullets = spec.get("bullets", [])

    # Title placeholder
    if title_text and slide.shapes.title is not None:
        slide.shapes.title.text = title_text

    # Classify remaining placeholders
    body_ph = None
    subtitle_ph = None
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        tname = str(t)
        if ph == slide.shapes.title:
            continue
        if "SUBTITLE" in tname and subtitle_ph is None:
            subtitle_ph = ph
        elif any(k in tname for k in ("BODY", "OBJECT", "CONTENT", "TEXT")) and body_ph is None:
            body_ph = ph

    if subtitle and subtitle_ph is not None:
        subtitle_ph.text = subtitle + cite_suffix(cites if not bullets else [])
    elif subtitle and body_ph is not None:
        body_ph.text = subtitle

    if bullets and body_ph is not None:
        tf = body_ph.text_frame
        tf.text = bullets[0] + (cite_suffix(cites) if len(bullets) == 1 else "")
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
        # citation line if multiple bullets
        if len(bullets) > 1 and cites:
            p = tf.add_paragraph()
            p.text = "Sources:" + cite_suffix(cites)

    notes = spec.get("notes")
    sources_note = ("Sources: " + ", ".join(cites)) if cites else ""
    note_text = "\n".join(x for x in (notes, sources_note) if x)
    if note_text:
        slide.notes_slide.notes_text_frame.text = note_text


def build(spec, template_path, out_override=None):
    prs = load_template(template_path)
    if spec.get("clear_template_slides", True):
        clear_slides(prs)
    lmap = layout_map(prs)
    for s in spec.get("slides", []):
        layout = pick_layout(lmap, s.get("layout", "Title and Content - White"))
        slide = prs.slides.add_slide(layout)
        fill_slide(slide, s)
    out = out_override or spec.get("output", "deck.pptx")
    if not os.path.isabs(out):
        out = os.path.join(os.getcwd(), out)
    prs.save(out)
    return out, len(spec.get("slides", []))


def main():
    if len(sys.argv) >= 2 and sys.argv[1] == "list-layouts":
        tpl = sys.argv[2] if len(sys.argv) > 2 else DEFAULT_TEMPLATE
        prs = load_template(tpl)
        for name in sorted(layout_map(prs)):
            print(name)
        return
    ap = argparse.ArgumentParser()
    ap.add_argument("spec")
    ap.add_argument("--template", default=DEFAULT_TEMPLATE)
    ap.add_argument("--out", default=None)
    args = ap.parse_args()
    spec = json.load(open(args.spec))
    out, n = build(spec, args.template, args.out)
    print(f"Wrote {out} ({n} slides) from {os.path.basename(args.template)}")


if __name__ == "__main__":
    main()
